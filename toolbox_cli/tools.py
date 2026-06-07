import os
import sys
import subprocess
import shutil
import re
from .base import (
    BaseConverter,
    get_matching_files,
    get_unique_output_path,
    select_files_to_convert,
    print_info,
    print_success,
    print_error,
    print_summary,
    HAS_RICH
)
from .i18n import get_text

# ==========================================
# 1. Image to ICO Converter
# ==========================================
class ImageToIcoConverter(BaseConverter):
    @property
    def id(self) -> int:
        return 10

    @property
    def name(self) -> str:
        return get_text("t_img_to_ico")

    @property
    def input_extensions(self) -> list:
        return [".png", ".jpg", ".jpeg", ".svg"]

    @property
    def output_extension_desc(self) -> str:
        return "ICO (16x16 to 256x256)"

    def convert(self, files: list) -> tuple:
        selected_files = select_files_to_convert(files, self.name)
        if not selected_files:
            return 0, 0

        try:
            from PIL import Image
        except ImportError:
            print_error("Pillow library is missing. Run install.bat first.")
            return 0, len(selected_files)

        success_count = 0
        fail_count = 0

        for filename in selected_files:
            try:
                output_path = get_unique_output_path(filename, ".ico")
                print_info(get_text("status_converting", input=filename, output=os.path.basename(output_path)))
                
                img = Image.open(filename)
                
                # Create a square canvas to avoid distortion
                width, height = img.size
                max_dim = max(width, height)
                square_img = Image.new("RGBA", (max_dim, max_dim), (0, 0, 0, 0))
                
                # Center the original image
                offset_x = (max_dim - width) // 2
                offset_y = (max_dim - height) // 2
                square_img.paste(img, (offset_x, offset_y))
                
                # Save as ICO with standard sizes
                icon_sizes = [(16, 16), (32, 32), (48, 48), (64, 64), (128, 128), (256, 256)]
                square_img.save(output_path, format="ICO", sizes=icon_sizes)
                
                print_success(get_text("status_converted", output=os.path.basename(output_path)))
                success_count += 1
            except Exception as e:
                print_error(get_text("err_general", file=filename, e=e))
                fail_count += 1

        return success_count, fail_count


# ==========================================
# 2. PDF to Image Converter
# ==========================================
class PdfToImageConverter(BaseConverter):
    @property
    def id(self) -> int:
        return 11

    @property
    def name(self) -> str:
        return get_text("t_pdf_to_img")

    @property
    def input_extensions(self) -> list:
        return [".pdf"]

    @property
    def output_extension_desc(self) -> str:
        return "PNG"

    def convert(self, files: list) -> tuple:
        selected_files = select_files_to_convert(files, self.name)
        if not selected_files:
            return 0, 0

        try:
            import fitz  # PyMuPDF
        except ImportError:
            print_error("PyMuPDF (fitz) library is missing. Run install.bat first.")
            return 0, len(selected_files)

        success_count = 0
        fail_count = 0

        for filename in selected_files:
            try:
                print_info(f"Extracting pages from PDF: {filename}...")
                doc = fitz.open(filename)
                name, _ = os.path.splitext(filename)
                
                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    # Use 300 DPI high resolution zoom matrix
                    zoom = 300 / 72
                    matrix = fitz.Matrix(zoom, zoom)
                    pix = page.get_pixmap(matrix=matrix)
                    
                    page_filename = f"{name}_page_{page_num + 1}.png"
                    unique_page_path = get_unique_output_path(page_filename, ".png")
                    
                    pix.save(unique_page_path)
                    print_success(get_text("status_converted", output=os.path.basename(unique_page_path)))
                
                doc.close()
                success_count += 1
            except Exception as e:
                print_error(get_text("err_general", file=filename, e=e))
                fail_count += 1

        return success_count, fail_count


# ==========================================
# 3. PDF to Markdown Converter
# ==========================================
class PdfToMarkdownConverter(BaseConverter):
    @property
    def id(self) -> int:
        return 5

    @property
    def name(self) -> str:
        return get_text("t_pdf_to_md")

    @property
    def input_extensions(self) -> list:
        return [".pdf"]

    @property
    def output_extension_desc(self) -> str:
        return "MD"

    def convert(self, files: list) -> tuple:
        selected_files = select_files_to_convert(files, self.name)
        if not selected_files:
            return 0, 0

        try:
            import fitz
        except ImportError:
            print_error("PyMuPDF (fitz) library is missing. Run install.bat first.")
            return 0, len(selected_files)

        success_count = 0
        fail_count = 0

        for filename in selected_files:
            try:
                output_path = get_unique_output_path(filename, ".md")
                print_info(get_text("status_converting", input=filename, output=os.path.basename(output_path)))
                
                doc = fitz.open(filename)
                markdown_lines = []
                
                # Add titles and metadata structure
                markdown_lines.append(f"# {os.path.splitext(filename)[0]}")
                markdown_lines.append("\n---\n")
                
                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    text = page.get_text("text")
                    
                    markdown_lines.append(f"\n## Page {page_num + 1}\n")
                    
                    # Heuristics for clean structure formatting
                    lines = text.split("\n")
                    for line in lines:
                        cleaned = line.strip()
                        if not cleaned:
                            continue
                        
                        # Detect potential headings (capitalized, short)
                        if len(cleaned) < 60 and cleaned.isupper() and not cleaned.endswith((".", ",", ";")):
                            markdown_lines.append(f"\n### {cleaned}\n")
                        # Detect list items
                        elif cleaned.startswith(("-", "*", "+", "•")) or (cleaned[0].isdigit() and cleaned.split(".")[0].isdigit()):
                            markdown_lines.append(f"\n* {cleaned.lstrip('-*+•0123456789. ')}")
                        else:
                            markdown_lines.append(cleaned)
                    
                    markdown_lines.append("\n")
                
                doc.close()
                
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write("\n".join(markdown_lines))
                    
                print_success(get_text("status_converted", output=os.path.basename(output_path)))
                success_count += 1
            except Exception as e:
                print_error(get_text("err_general", file=filename, e=e))
                fail_count += 1

        return success_count, fail_count


# ==========================================
# 4. PowerPoint to Markdown Converter
# ==========================================
class PptToMarkdownConverter(BaseConverter):
    @property
    def id(self) -> int:
        return 4

    @property
    def name(self) -> str:
        return get_text("t_ppt_to_md")

    @property
    def input_extensions(self) -> list:
        return [".pptx", ".ppt"]

    @property
    def output_extension_desc(self) -> str:
        return "MD"

    def convert(self, files: list) -> tuple:
        selected_files = select_files_to_convert(files, self.name)
        if not selected_files:
            return 0, 0

        try:
            from pptx import Presentation
        except ImportError:
            print_error("python-pptx library is missing. Run install.bat first.")
            return 0, len(selected_files)

        success_count = 0
        fail_count = 0

        for filename in selected_files:
            # Check legacy vs modern ppt format
            _, ext = os.path.splitext(filename)
            pptx_filename = filename
            temp_created = False
            
            # Convert legacy .ppt to .pptx via Word COM if on Windows
            if ext.lower() == ".ppt":
                if sys.platform.startswith('win'):
                    try:
                        import win32com.client
                        print_info("Converting legacy .ppt to .pptx using PowerPoint COM...")
                        powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
                        abs_in = os.path.abspath(filename)
                        temp_pptx = os.path.abspath("temp_conversion_powerpoint.pptx")
                        
                        pres = powerpoint.Presentations.Open(abs_in, WithWindow=False)
                        pres.SaveAs(temp_pptx, 24)  # 24 = ppSaveAsOpenXMLPresentation
                        pres.Close()
                        powerpoint.Quit()
                        
                        pptx_filename = "temp_conversion_powerpoint.pptx"
                        temp_created = True
                    except Exception as e:
                        print_error(f"Could not convert legacy PowerPoint file: {e}")
                        fail_count += 1
                        continue
                else:
                    print_error("Converting legacy .ppt requires Microsoft PowerPoint on Windows.")
                    fail_count += 1
                    continue

            try:
                output_path = get_unique_output_path(filename, ".md")
                print_info(get_text("status_converting", input=filename, output=os.path.basename(output_path)))
                
                prs = Presentation(pptx_filename)
                markdown_lines = []
                markdown_lines.append(f"# Presentation: {os.path.splitext(filename)[0]}")
                markdown_lines.append("\n---\n")
                
                for slide_num, slide in enumerate(prs.slides, 1):
                    markdown_lines.append(f"\n## Slide {slide_num}\n")
                    
                    # Extract text from slide shapes
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text = shape.text.strip()
                            # If shape is slide title
                            if shape == slide.shapes[0] or "Title" in shape.name:
                                markdown_lines.append(f"### {text}\n")
                            else:
                                markdown_lines.append(text)
                    
                    markdown_lines.append("\n---\n")
                
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write("\n".join(markdown_lines))
                    
                print_success(get_text("status_converted", output=os.path.basename(output_path)))
                success_count += 1
            except Exception as e:
                print_error(get_text("err_general", file=filename, e=e))
                fail_count += 1
            finally:
                if temp_created and os.path.exists("temp_conversion_powerpoint.pptx"):
                    try:
                        os.remove("temp_conversion_powerpoint.pptx")
                    except Exception:
                        pass

        return success_count, fail_count


# ==========================================
# 5. Word to Markdown Converter
# ==========================================
class DocxToMarkdownConverter(BaseConverter):
    @property
    def id(self) -> int:
        return 3

    @property
    def name(self) -> str:
        return get_text("t_docx_to_md")

    @property
    def input_extensions(self) -> list:
        return [".docx", ".doc"]

    @property
    def output_extension_desc(self) -> str:
        return "MD"

    def convert(self, files: list) -> tuple:
        selected_files = select_files_to_convert(files, self.name)
        if not selected_files:
            return 0, 0

        try:
            from docx import Document
        except ImportError:
            print_error("python-docx library is missing. Run install.bat first.")
            return 0, len(selected_files)

        success_count = 0
        fail_count = 0

        for filename in selected_files:
            _, ext = os.path.splitext(filename)
            docx_filename = filename
            temp_created = False
            
            # Convert legacy doc to docx using Word COM
            if ext.lower() == ".doc":
                if sys.platform.startswith('win'):
                    try:
                        import win32com.client
                        print_info("Converting legacy .doc to .docx using Word COM...")
                        word = win32com.client.DispatchEx("Word.Application")
                        abs_in = os.path.abspath(filename)
                        temp_docx = os.path.abspath("temp_conversion_word.docx")
                        
                        doc = word.Documents.Open(abs_in)
                        doc.SaveAs2(temp_docx, FileFormat=16)  # 16 = wdFormatXMLDocument
                        doc.Close()
                        word.Quit()
                        
                        docx_filename = "temp_conversion_word.docx"
                        temp_created = True
                    except Exception as e:
                        print_error(f"Could not convert legacy Word file: {e}")
                        fail_count += 1
                        continue
                else:
                    print_error("Converting legacy .doc requires Microsoft Word on Windows.")
                    fail_count += 1
                    continue

            try:
                output_path = get_unique_output_path(filename, ".md")
                print_info(get_text("status_converting", input=filename, output=os.path.basename(output_path)))
                
                doc = Document(docx_filename)
                markdown_lines = []
                
                # Title metadata
                markdown_lines.append(f"# Document: {os.path.splitext(filename)[0]}")
                markdown_lines.append("\n---\n")
                
                for p in doc.paragraphs:
                    style = p.style.name.lower()
                    text = p.text.strip()
                    if not text:
                        continue
                    
                    if "heading 1" in style:
                        markdown_lines.append(f"\n# {text}\n")
                    elif "heading 2" in style:
                        markdown_lines.append(f"\n## {text}\n")
                    elif "heading 3" in style:
                        markdown_lines.append(f"\n### {text}\n")
                    elif "list bullet" in style:
                        markdown_lines.append(f"* {text}")
                    elif "list number" in style:
                        markdown_lines.append(f"1. {text}")
                    else:
                        markdown_lines.append(text)
                
                # Parse tables in document
                for table in doc.tables:
                    markdown_lines.append("\n")
                    headers = [cell.text.strip() for cell in table.rows[0].cells]
                    markdown_lines.append("| " + " | ".join(headers) + " |")
                    markdown_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
                    
                    for row in table.rows[1:]:
                        row_cells = [cell.text.strip().replace("\n", "<br>") for cell in row.cells]
                        markdown_lines.append("| " + " | ".join(row_cells) + " |")
                    markdown_lines.append("\n")

                with open(output_path, "w", encoding="utf-8") as f:
                    f.write("\n".join(markdown_lines))
                    
                print_success(get_text("status_converted", output=os.path.basename(output_path)))
                success_count += 1
            except Exception as e:
                print_error(get_text("err_general", file=filename, e=e))
                fail_count += 1
            finally:
                if temp_created and os.path.exists("temp_conversion_word.docx"):
                    try:
                        os.remove("temp_conversion_word.docx")
                    except Exception:
                        pass

        return success_count, fail_count


# ==========================================
# 6. Excel to CSV Converter
# ==========================================
class ExcelToCsvConverter(BaseConverter):
    @property
    def id(self) -> int:
        return 7

    @property
    def name(self) -> str:
        return get_text("t_xls_to_csv")

    @property
    def input_extensions(self) -> list:
        return [".xlsx", ".xls"]

    @property
    def output_extension_desc(self) -> str:
        return "CSV"

    def convert(self, files: list) -> tuple:
        selected_files = select_files_to_convert(files, self.name)
        if not selected_files:
            return 0, 0

        try:
            import pandas as pd
        except ImportError:
            print_error("pandas library is missing. Run install.bat first.")
            return 0, len(selected_files)

        success_count = 0
        fail_count = 0

        for filename in selected_files:
            try:
                print_info(f"Reading workbook sheets: {filename}...")
                excel_file = pd.ExcelFile(filename)
                name, _ = os.path.splitext(filename)
                
                for sheet_name in excel_file.sheet_names:
                    df = excel_file.parse(sheet_name)
                    # Clean sheet name for filename compatibility
                    clean_sheet = re.sub(r'[\\/*?:\[\]]', "_", sheet_name)
                    csv_filename = f"{name}_{clean_sheet}.csv"
                    output_path = get_unique_output_path(csv_filename, ".csv")
                    
                    df.to_csv(output_path, index=False, encoding="utf-8")
                    print_success(get_text("status_converted", output=os.path.basename(output_path)))
                    
                success_count += 1
            except Exception as e:
                print_error(get_text("err_general", file=filename, e=e))
                fail_count += 1

        return success_count, fail_count


# ==========================================
# 7. CSV to Excel Converter
# ==========================================
class CsvToExcelConverter(BaseConverter):
    @property
    def id(self) -> int:
        return 8

    @property
    def name(self) -> str:
        return get_text("t_csv_to_xls")

    @property
    def input_extensions(self) -> list:
        return [".csv"]

    @property
    def output_extension_desc(self) -> str:
        return "XLSX"

    def convert(self, files: list) -> tuple:
        selected_files = select_files_to_convert(files, self.name)
        if not selected_files:
            return 0, 0

        try:
            import pandas as pd
        except ImportError:
            print_error("pandas library is missing. Run install.bat first.")
            return 0, len(selected_files)

        success_count = 0
        fail_count = 0

        for filename in selected_files:
            try:
                output_path = get_unique_output_path(filename, ".xlsx")
                print_info(get_text("status_converting", input=filename, output=os.path.basename(output_path)))
                
                df = pd.read_csv(filename, encoding="utf-8")
                
                # Write to Excel format using openpyxl engine
                with pd.ExcelWriter(output_path, engine="openpyxl") as writer:
                    df.to_excel(writer, sheet_name="Data", index=False)
                    
                print_success(get_text("status_converted", output=os.path.basename(output_path)))
                success_count += 1
            except Exception as e:
                print_error(get_text("err_general", file=filename, e=e))
                fail_count += 1

        return success_count, fail_count


# ==========================================
# 8. Markdown to HTML Converter
# ==========================================
class MarkdownToHtmlConverter(BaseConverter):
    @property
    def id(self) -> int:
        return 6

    @property
    def name(self) -> str:
        return get_text("t_md_to_html")

    @property
    def input_extensions(self) -> list:
        return [".md"]

    @property
    def output_extension_desc(self) -> str:
        return "HTML"

    def convert(self, files: list) -> tuple:
        selected_files = select_files_to_convert(files, self.name)
        if not selected_files:
            return 0, 0

        try:
            import markdown
        except ImportError:
            print_error("markdown library is missing. Run install.bat first.")
            return 0, len(selected_files)

        success_count = 0
        fail_count = 0

        # Premium CSS stylesheet with responsive layout, code copy button, and auto dark-mode
        premium_css = """
        :root {
            --bg-color: #f8fafc;
            --text-color: #1e293b;
            --card-bg: #ffffff;
            --primary-color: #2563eb;
            --code-bg: #0f172a;
            --border-color: #e2e8f0;
            --header-bg: #ffffff;
        }
        @media (prefers-color-scheme: dark) {
            :root {
                --bg-color: #0f172a;
                --text-color: #f1f5f9;
                --card-bg: #1e293b;
                --primary-color: #3b82f6;
                --code-bg: #020617;
                --border-color: #334155;
                --header-bg: #1e293b;
            }
        }
        body {
            font-family: 'Inter', -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
            background-color: var(--bg-color);
            color: var(--text-color);
            line-height: 1.7;
            margin: 0;
            padding: 0;
        }
        .container {
            max-width: 800px;
            margin: 40px auto;
            padding: 40px;
            background-color: var(--card-bg);
            border-radius: 12px;
            box-shadow: 0 10px 15px -3px rgba(0, 0, 0, 0.1), 0 4px 6px -4px rgba(0, 0, 0, 0.1);
            border: 1px solid var(--border-color);
        }
        h1, h2, h3 {
            color: var(--primary-color);
            margin-top: 1.5em;
        }
        h1 {
            font-size: 2.25rem;
            border-bottom: 2px solid var(--border-color);
            padding-bottom: 0.3em;
        }
        a {
            color: var(--primary-color);
            text-decoration: none;
        }
        a:hover {
            text-decoration: underline;
        }
        pre {
            background-color: var(--code-bg);
            color: #f8f8f2;
            padding: 1.25em;
            border-radius: 8px;
            overflow-x: auto;
            position: relative;
        }
        code {
            font-family: 'Fira Code', 'Courier New', Courier, monospace;
            font-size: 0.9em;
        }
        .copy-btn {
            position: absolute;
            top: 10px;
            right: 10px;
            background-color: var(--primary-color);
            color: white;
            border: none;
            border-radius: 4px;
            padding: 5px 10px;
            font-size: 0.75rem;
            cursor: pointer;
            opacity: 0.6;
            transition: opacity 0.2s;
        }
        .copy-btn:hover {
            opacity: 1;
        }
        table {
            width: 100%;
            border-collapse: collapse;
            margin: 1.5em 0;
        }
        th, td {
            border: 1px solid var(--border-color);
            padding: 12px;
            text-align: left;
        }
        th {
            background-color: var(--primary-color);
            color: white;
        }
        tr:nth-child(even) {
            background-color: var(--bg-color);
        }
        """

        premium_js = """
        document.querySelectorAll('pre').forEach((preBlock) => {
            const button = document.createElement('button');
            button.className = 'copy-btn';
            button.textContent = 'Copy';
            preBlock.appendChild(button);
            
            button.addEventListener('click', () => {
                const code = preBlock.querySelector('code');
                const text = code ? code.innerText : preBlock.innerText.replace('Copy', '').trim();
                
                navigator.clipboard.writeText(text).then(() => {
                    button.textContent = 'Copied!';
                    setTimeout(() => {
                        button.textContent = 'Copy';
                    }, 2000);
                });
            });
        });
        """

        for filename in selected_files:
            try:
                output_path = get_unique_output_path(filename, ".html")
                print_info(get_text("status_converting", input=filename, output=os.path.basename(output_path)))
                
                with open(filename, "r", encoding="utf-8") as f:
                    md_text = f.read()
                    
                # Convert Markdown text to HTML body
                html_body = markdown.markdown(md_text, extensions=["extra", "tables", "codehilite"])
                
                page_title = os.path.splitext(filename)[0]
                
                # Formulate structural premium html markup
                full_html = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{page_title}</title>
    <meta name="description" content="Generated from {filename} by Toolbox CLI.">
    <link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&family=Fira+Code&display=swap" rel="stylesheet">
    <style>
        {premium_css}
    </style>
</head>
<body>
    <div class="container">
        {html_body}
    </div>
    <script>
        {premium_js}
    </script>
</body>
</html>
"""
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write(full_html)
                    
                print_success(get_text("status_converted", output=os.path.basename(output_path)))
                success_count += 1
            except Exception as e:
                print_error(get_text("err_general", file=filename, e=e))
                fail_count += 1

        return success_count, fail_count


# ==========================================
# 9. General Image Format Converter
# ==========================================
class ImageConverter(BaseConverter):
    @property
    def id(self) -> int:
        return 9

    @property
    def name(self) -> str:
        return get_text("t_img_conv")

    @property
    def input_extensions(self) -> list:
        return [".png", ".jpg", ".jpeg", ".webp", ".bmp"]

    @property
    def output_extension_desc(self) -> str:
        return "PNG, JPG, WEBP, BMP"

    def convert(self, files: list) -> tuple:
        selected_files = select_files_to_convert(files, self.name)
        if not selected_files:
            return 0, 0

        # Ask target format
        print()
        if HAS_RICH:
            target_ext = Prompt.ask(get_text("img_conv_prompt"), choices=["png", "jpg", "webp", "bmp"]).strip().lower()
        else:
            target_ext = input(f"{get_text('img_conv_prompt')}: ").strip().lower()

        if target_ext not in ["png", "jpg", "webp", "bmp"]:
            print_error("Invalid format choice.")
            return 0, len(selected_files)

        target_ext = f".{target_ext}"
        
        try:
            from PIL import Image
        except ImportError:
            print_error("Pillow library is missing. Run install.bat first.")
            return 0, len(selected_files)

        success_count = 0
        fail_count = 0

        for filename in selected_files:
            try:
                output_path = get_unique_output_path(filename, target_ext)
                print_info(get_text("status_converting", input=filename, output=os.path.basename(output_path)))
                
                img = Image.open(filename)
                
                # Handle color profiles and alpha channel for JPG conversion
                if target_ext == ".jpg" and img.mode in ("RGBA", "LA", "P"):
                    background = Image.new("RGB", img.size, (255, 255, 255))
                    # Mask alpha paste
                    if img.mode == "RGBA":
                        background.paste(img, mask=img.split()[3])
                    else:
                        background.paste(img.convert("RGBA"), mask=img.convert("RGBA").split()[3])
                    img = background
                elif img.mode != "RGBA" and target_ext in (".png", ".webp"):
                    img = img.convert("RGBA")
                    
                img.save(output_path)
                print_success(get_text("status_converted", output=os.path.basename(output_path)))
                success_count += 1
            except Exception as e:
                print_error(get_text("err_general", file=filename, e=e))
                fail_count += 1

        return success_count, fail_count


# ==========================================
# 10. Intelligent Image Compression
# ==========================================
class ImageCompressionConverter(BaseConverter):
    @property
    def id(self) -> int:
        return 12

    @property
    def name(self) -> str:
        return get_text("t_img_comp")

    @property
    def input_extensions(self) -> list:
        return [".png", ".jpg", ".jpeg", ".webp"]

    @property
    def output_extension_desc(self) -> str:
        return "Compressed Images"

    def convert(self, files: list) -> tuple:
        selected_files = select_files_to_convert(files, self.name)
        if not selected_files:
            return 0, 0

        print()
        if HAS_RICH:
            mode = Prompt.ask(get_text("img_comp_prompt"), choices=["1", "2", "3"]).strip()
        else:
            mode = input(f"{get_text('img_comp_prompt')} (1-3): ").strip()

        quality = 85
        if mode == "1":
            quality = 100  # Lossless
        elif mode == "2":
            quality = 65   # High compression
        elif mode == "3":
            try:
                if HAS_RICH:
                    quality = IntPrompt.ask(get_text("img_quality_prompt"))
                else:
                    quality = int(input(f"{get_text('img_quality_prompt')}: ").strip())
            except ValueError:
                quality = 85

        try:
            from PIL import Image
        except ImportError:
            print_error("Pillow library is missing. Run install.bat first.")
            return 0, len(selected_files)

        success_count = 0
        fail_count = 0

        for filename in selected_files:
            try:
                name, ext = os.path.splitext(filename)
                output_path = get_unique_output_path(f"{name}_compressed", ext.lower())
                print_info(get_text("status_converting", input=filename, output=os.path.basename(output_path)))
                
                img = Image.open(filename)
                
                # Check extension and save using specific compression engines
                if ext.lower() in (".jpg", ".jpeg"):
                    img.save(output_path, "JPEG", quality=quality, optimize=True)
                elif ext.lower() == ".png":
                    if quality == 100:
                        img.save(output_path, "PNG", optimize=True)
                    else:
                        # Convert to adaptive quantized palette format (lossy but very light)
                        quantized = img.quantize(colors=256)
                        quantized.save(output_path, "PNG", optimize=True)
                elif ext.lower() == ".webp":
                    img.save(output_path, "WEBP", quality=quality, lossless=(quality == 100))
                
                orig_size = os.path.getsize(filename)
                new_size = os.path.getsize(output_path)
                pct = (1 - (new_size / orig_size)) * 100
                
                print_success(f"{get_text('status_converted', output=os.path.basename(output_path))} ({pct:.1f}% space saved)")
                success_count += 1
            except Exception as e:
                print_error(get_text("err_general", file=filename, e=e))
                fail_count += 1

        return success_count, fail_count


# ==========================================
# 11. Intelligent Video Compression (FFmpeg)
# ==========================================
class VideoCompressionConverter(BaseConverter):
    @property
    def id(self) -> int:
        return 13

    @property
    def name(self) -> str:
        return get_text("t_vid_comp")

    @property
    def input_extensions(self) -> list:
        return [".mp4", ".mkv", ".avi", ".mov", ".wmv", ".flv", ".3gp", ".webm"]

    @property
    def output_extension_desc(self) -> str:
        return "MP4 (H.264 or H.265)"

    def convert(self, files: list) -> tuple:
        # Check if FFmpeg is in system PATH
        if not shutil.which("ffmpeg"):
            print_error(get_text("vid_err_ffmpeg"))
            return 0, len(files)

        selected_files = select_files_to_convert(files, self.name)
        if not selected_files:
            return 0, 0

        print()
        if HAS_RICH:
            codec_choice = Prompt.ask(get_text("vid_codec_prompt"), choices=["1", "2"]).strip()
        else:
            codec_choice = input(f"{get_text('vid_codec_prompt')} (1-2): ").strip()

        codec = "libx265" if codec_choice == "1" else "libx264"

        if HAS_RICH:
            crf_choice = Prompt.ask(get_text("vid_crf_prompt"), choices=["1", "2", "3"]).strip()
        else:
            crf_choice = input(f"{get_text('vid_crf_prompt')} (1-3): ").strip()

        # Set Constant Rate Factor (CRF) for H.264 / H.265 quality targeting
        if crf_choice == "1":
            crf = "20"
        elif crf_choice == "2":
            crf = "24"
        else:
            crf = "28"

        success_count = 0
        fail_count = 0

        for filename in selected_files:
            try:
                name, _ = os.path.splitext(filename)
                output_path = get_unique_output_path(f"{name}_compressed", ".mp4")
                print_info(get_text("status_converting", input=filename, output=os.path.basename(output_path)))
                
                # Run local FFmpeg subprocess cleanly
                cmd = [
                    "ffmpeg", "-y",
                    "-i", filename,
                    "-vcodec", codec,
                    "-crf", crf,
                    "-acodec", "aac",
                    "-b:a", "128k",
                    output_path
                ]
                
                # Prevent popup console windows in Windows environments
                creationflags = subprocess.CREATE_NO_WINDOW if sys.platform.startswith('win') else 0
                
                subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, creationflags=creationflags)
                
                orig_size = os.path.getsize(filename)
                new_size = os.path.getsize(output_path)
                pct = (1 - (new_size / orig_size)) * 100
                
                print_success(f"{get_text('status_converted', output=os.path.basename(output_path))} ({pct:.1f}% space saved)")
                success_count += 1
            except Exception as e:
                print_error(get_text("err_general", file=filename, e=e))
                fail_count += 1

        return success_count, fail_count


# ==========================================
# 12. Word to PDF (Native COM / LibreOffice)
# ==========================================
class WordToPdfConverter(BaseConverter):
    @property
    def id(self) -> int:
        return 1

    @property
    def name(self) -> str:
        return get_text("t_docx_to_pdf")

    @property
    def input_extensions(self) -> list:
        return [".docx", ".doc"]

    @property
    def output_extension_desc(self) -> str:
        return "PDF"

    def convert(self, files: list) -> tuple:
        selected_files = select_files_to_convert(files, self.name)
        if not selected_files:
            return 0, 0

        success_count = 0
        fail_count = 0

        if sys.platform.startswith('win'):
            try:
                import win32com.client
            except ImportError:
                print_error(get_text("err_no_pywin32"))
                return 0, len(selected_files)

            print_info(get_text("status_com_start"))
            try:
                # DispatchEx ensures COM isolation in Windows
                word = win32com.client.DispatchEx("Word.Application")
            except Exception as e:
                print_error(get_text("err_com_init", e=e))
                return 0, len(selected_files)

            for filename in selected_files:
                filepath = os.path.abspath(filename)
                pdf_path = get_unique_output_path(filename, ".pdf")
                
                # Resolve short directory paths to bypass emojis bugs inside win32com client
                dir_path = os.path.dirname(filepath)
                import ctypes
                buf = ctypes.create_unicode_buffer(500)
                ctypes.windll.kernel32.GetShortPathNameW(dir_path, buf, 500)
                short_dir = buf.value
                
                if short_dir:
                    filepath_short = os.path.join(short_dir, os.path.basename(filepath))
                    pdf_path_short = os.path.join(short_dir, os.path.basename(pdf_path))
                else:
                    filepath_short = filepath
                    pdf_path_short = pdf_path
                    
                try:
                    print_info(get_text("status_converting", input=filename, output=os.path.basename(pdf_path)))
                    doc = word.Documents.Open(filepath_short)
                    doc.SaveAs2(pdf_path_short, FileFormat=17) # 17 = wdFormatPDF
                    doc.Close()
                    print_success(get_text("status_converted", output=os.path.basename(pdf_path)))
                    success_count += 1
                except Exception as e:
                    print_error(get_text("err_general", file=filename, e=e))
                    fail_count += 1
                    
            try:
                word.Quit()
            except Exception:
                pass
        else:
            # Linux / macOS fallback using LibreOffice headless command
            if not shutil.which("libreoffice"):
                print_error("LibreOffice binary not found. Headless docx-pdf conversion requires LibreOffice.")
                return 0, len(selected_files)
                
            for filename in selected_files:
                try:
                    pdf_path = get_unique_output_path(filename, ".pdf")
                    print_info(get_text("status_converting", input=filename, output=os.path.basename(pdf_path)))
                    
                    cmd = ["libreoffice", "--headless", "--convert-to", "pdf", filename]
                    subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
                    
                    # Move output if renamed
                    generated_name = os.path.splitext(filename)[0] + ".pdf"
                    if generated_name != pdf_path and os.path.exists(generated_name):
                        shutil.move(generated_name, pdf_path)
                        
                    print_success(get_text("status_converted", output=os.path.basename(pdf_path)))
                    success_count += 1
                except Exception as e:
                    print_error(get_text("err_general", file=filename, e=e))
                    fail_count += 1

        return success_count, fail_count


# ==========================================
# 13. PDF to Word (Native COM / LibreOffice)
# ==========================================
class PdfToWordConverter(BaseConverter):
    @property
    def id(self) -> int:
        return 2

    @property
    def name(self) -> str:
        return get_text("t_pdf_to_docx")

    @property
    def input_extensions(self) -> list:
        return [".pdf"]

    @property
    def output_extension_desc(self) -> str:
        return "DOCX"

    def convert(self, files: list) -> tuple:
        selected_files = select_files_to_convert(files, self.name)
        if not selected_files:
            return 0, 0

        success_count = 0
        fail_count = 0

        if sys.platform.startswith('win'):
            # Bypass pdf import warning dialogs in Microsoft Word Options registry
            import winreg
            for v in ["15.0", "16.0", "17.0"]:
                try:
                    key_path = rf"Software\Microsoft\Office\{v}\Word\Options"
                    key = winreg.OpenKey(winreg.HKEY_CURRENT_USER, key_path, 0, winreg.KEY_SET_VALUE)
                    winreg.SetValueEx(key, "DisableConvertPdfWarning", 0, winreg.REG_DWORD, 1)
                    winreg.CloseKey(key)
                except Exception:
                    try:
                        # Create if key does not exist
                        key_path = rf"Software\Microsoft\Office\{v}\Word\Options"
                        key = winreg.CreateKey(winreg.HKEY_CURRENT_USER, key_path)
                        winreg.SetValueEx(key, "DisableConvertPdfWarning", 0, winreg.REG_DWORD, 1)
                        winreg.CloseKey(key)
                    except Exception:
                        pass
                        
            try:
                import win32com.client
            except ImportError:
                print_error(get_text("err_no_pywin32"))
                return 0, len(selected_files)

            print_info(get_text("status_com_start"))
            try:
                word = win32com.client.DispatchEx("Word.Application")
            except Exception as e:
                print_error(get_text("err_com_init", e=e))
                return 0, len(selected_files)

            for filename in selected_files:
                filepath = os.path.abspath(filename)
                docx_path = get_unique_output_path(filename, ".docx")
                
                # Resolve short directory paths to bypass emojis bugs inside win32com client
                dir_path = os.path.dirname(filepath)
                import ctypes
                buf = ctypes.create_unicode_buffer(500)
                ctypes.windll.kernel32.GetShortPathNameW(dir_path, buf, 500)
                short_dir = buf.value
                
                if short_dir:
                    filepath_short = os.path.join(short_dir, os.path.basename(filepath))
                    docx_path_short = os.path.join(short_dir, os.path.basename(docx_path))
                else:
                    filepath_short = filepath
                    docx_path_short = docx_path
                    
                try:
                    print_info(get_text("status_converting", input=filename, output=os.path.basename(docx_path)))
                    doc = word.Documents.Open(filepath_short, ConfirmConversions=False)
                    doc.SaveAs2(docx_path_short, FileFormat=16) # 16 = wdFormatXMLDocument
                    doc.Close()
                    print_success(get_text("status_converted", output=os.path.basename(docx_path)))
                    success_count += 1
                except Exception as e:
                    print_error(get_text("err_general", file=filename, e=e))
                    fail_count += 1
                    
            try:
                word.Quit()
            except Exception:
                pass
        else:
            # Fallback for Linux / macOS - LibreOffice conversion from PDF to Word is poor, 
            # so we inform the user or do a pdf-docx import if LibreOffice handles it
            if not shutil.which("libreoffice"):
                print_error("LibreOffice binary not found. PDF-docx conversion requires LibreOffice on Unix.")
                return 0, len(selected_files)
                
            for filename in selected_files:
                try:
                    docx_path = get_unique_output_path(filename, ".docx")
                    print_info(get_text("status_converting", input=filename, output=os.path.basename(docx_path)))
                    
                    cmd = ["libreoffice", "--headless", "--convert-to", "docx", filename]
                    subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
                    
                    generated_name = os.path.splitext(filename)[0] + ".docx"
                    if generated_name != docx_path and os.path.exists(generated_name):
                        shutil.move(generated_name, docx_path)
                        
                    print_success(get_text("status_converted", output=os.path.basename(docx_path)))
                    success_count += 1
                except Exception as e:
                    print_error(get_text("err_general", file=filename, e=e))
                    fail_count += 1

        return success_count, fail_count


# ==========================================
# Strategy Registry
# ==========================================
CONVERTERS = [
    WordToPdfConverter(),
    PdfToWordConverter(),
    DocxToMarkdownConverter(),
    PptToMarkdownConverter(),
    PdfToMarkdownConverter(),
    MarkdownToHtmlConverter(),
    ExcelToCsvConverter(),
    CsvToExcelConverter(),
    ImageConverter(),
    ImageToIcoConverter(),
    PdfToImageConverter(),
    ImageCompressionConverter(),
    VideoCompressionConverter()
]

def get_converter_by_id(cid: int):
    """Find and return converter matching choice ID."""
    for c in CONVERTERS:
        if c.id == cid:
            return c
    return None
